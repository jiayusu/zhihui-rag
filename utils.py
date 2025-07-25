"""
- 扩展文件类型支持(DOCX、PPTX等)
- 提升文本提取准确率
- 添加文件元数据提取功能
- 实现更智能的文本分块策略
- 增加内容验证与清洗机制
- 添加并行文件处理能力
- 实现大文件处理进度跟踪
- 增加损坏文件处理机制
- 添加文件大小限制与验证
"""
from cgitb import text
import os
from typing import Dict, List, Optional, Tuple, Union

import PyPDF2
import pdfplumber
import markdown
import html2text
from docx import Document
from pptx import Presentation
import time
import concurrent.futures
from typing import Dict, List, Optional, Tuple, Union, Generator
import json
from tqdm import tqdm
import tiktoken
from bs4 import BeautifulSoup
import re

enc = tiktoken.get_encoding("cl100k_base")


class ReadFiles:
    """
    class to read files
    """

    def __init__(self, path: str) -> None:
        self._path = path
        self.file_list = self.get_files()

    def get_files(self):
        # args：dir_path，目标文件夹路径
        file_list = []
        for filepath, dirnames, filenames in os.walk(self._path):
            # os.walk 函数将递归遍历指定文件夹
            for filename in filenames:
                # 通过后缀名判断文件类型是否满足要求
                # 支持的文件类型列表
                supported_extensions = {'.md', '.txt', '.pdf', '.docx', '.pptx'}
                ext = os.path.splitext(filename)[1].lower()
                if ext in supported_extensions:
                    file_path = os.path.join(filepath, filename)
                    file_list.append(file_path)
        return file_list

    def get_content(self, path: str, max_token_len: int = 600, cover_content: int = 150, max_file_size: int = 100 * 1024 * 1024, max_workers: int = 4):
        """
        读取所有文件内容并分块，包含元数据提取和文件大小验证，支持并行处理
        
        Args:
            max_token_len: 最大令牌长度
            cover_content: 块重叠内容长度
            max_file_size: 最大文件大小限制(字节)
            max_workers: 并行处理的最大工作线程数
        
        Returns:
            包含分块内容和元数据的文档列表
        """
        docs = []
        
        # 使用线程池并行处理文件
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            # 提交所有文件处理任务
            futures = {
                executor.submit(
                    self._process_single_file,
                    file,
                    max_token_len,
                    cover_content,
                    max_file_size
                ): file for file in self.file_list
            }
            
            # 使用tqdm跟踪进度
            for future in tqdm(
                concurrent.futures.as_completed(futures),
                total=len(futures),
                desc="处理文件"
            ):
                file = futures[future]
                try:
                    result = future.result()
                    if result:
                        docs.extend(result)
                except Exception as e:
                    self.logger.error(f"处理文件 {file} 时出错: {str(e)}", exc_info=True)
                    continue
        
        return docs

    def _process_single_file(self, file: str, max_token_len: int, cover_content: int, max_file_size: int) -> Optional[List[Dict]]:
        """处理单个文件的辅助方法，用于并行处理"""
        try:
            # 文件大小验证
            file_size = os.path.getsize(file)
            if file_size > max_file_size:
                self.logger.warning(f"文件 {file} 大小超过限制 {max_file_size} bytes，已跳过")
                return None
            
            # 提取元数据
            file_stat = os.stat(file)
            metadata = {
                'file_path': file,
                'file_name': os.path.basename(file),
                'file_type': os.path.splitext(file)[1].lower(),
                'file_size': file_size,
                'created_time': time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(file_stat.st_ctime)),
                'modified_time': time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(file_stat.st_mtime)),
                'page_count': self._get_page_count(file)
            }
            
            # 读取文件内容
            content = self.read_file_content(file)
            
            # 内容验证与清洗
            cleaned_content = self.clean_and_validate_text(content)
            if not cleaned_content:
                self.logger.warning(f"文件 {file} 内容为空或清洗后为空，已跳过")
                return None
            
            # 智能分块
            chunk_content = self.get_chunk(
                cleaned_content, max_token_len=max_token_len, cover_content=cover_content)
            
            # 添加元数据到每个块
            file_docs = []
            for i, chunk in enumerate(chunk_content):
                doc = {
                    'content': chunk,
                    'metadata': {
                        **metadata,
                        'chunk_index': i,
                        'total_chunks': len(chunk_content),
                        'chunk_token_count': len(enc.encode(chunk))
                    }
                }
                file_docs.append(doc)
            
            return file_docs
        except Exception as e:
            self.logger.error(f"处理文件 {file} 时出错: {str(e)}", exc_info=True)
            return None

    def _get_page_count(self, file_path: str) -> Optional[int]:
        """获取文件页数(仅支持PDF和PPTX)"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            if ext == '.pdf':
                with pdfplumber.open(file_path) as pdf:
                    return len(pdf.pages)
            elif ext == '.pptx':
                prs = Presentation(file_path)
                return len(prs.slides)
            return None
        except Exception as e:
            self.logger.warning(f"获取文件页数失败 {file_path}: {str(e)}")
            return None

    @classmethod
    def clean_and_validate_text(cls, text: str) -> str:
        """内容验证与清洗"""
        if not isinstance(text, str):
            return ''
        
        # 移除多余空白字符
        text = re.sub(r'\s+', ' ', text).strip()
        
        # 移除特殊字符和控制字符
        text = re.sub(r'[\x00-\x1F\x7F]', '', text)
        
        # 验证文本长度
        if len(text) < 10:
            return ''
        
        return text

    @classmethod
    def get_chunk(cls,
                  text: str,
                  max_token_len: int = 600,
                  cover_content: int = 150,
                  separators: List[str] = ['\n\n', '\n', '. ', '! ', '? ', '; ', ', ', ' ']) -> List[str]:
        """
        智能文本分块，基于语义边界和令牌限制
        
        Args:
            text: 待分块文本
            max_token_len: 最大令牌长度
            cover_content: 块重叠内容长度
            separators: 分割符优先级列表
        
        Returns:
           分块后的文本列表
        """
        chunk_text = []
        token_len = max_token_len - cover_content
        current_chunk = ''
        current_tokens = 0
        
        # 按段落分割文本
        paragraphs = cls._split_text(text.strip(), separators)
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            para_tokens = len(enc.encode(para))
            
            # 如果段落本身超过最大长度，递归分块
            if para_tokens > max_token_len:
                sub_chunks = cls.get_chunk(
                    para,
                    max_token_len=max_token_len,
                    cover_content=cover_content,
                    separators=separators[1:]
                )
                for sub_chunk in sub_chunks:
                    chunk_text.append(sub_chunk)
                continue
            
            # 如果添加当前段落后超过限制，保存当前块并开始新块
            if current_tokens + para_tokens > token_len:
                chunk_text.append(current_chunk)
                # 新块包含重叠内容
                current_chunk = cls._get_cover_content(current_chunk, cover_content) + para
                current_tokens = len(enc.encode(current_chunk))
            else:
                current_chunk += para + '\n'
                current_tokens += para_tokens
        
        # 添加最后一个块
        if current_chunk.strip():
            chunk_text.append(current_chunk.strip())
        
        return chunk_text

    @classmethod
    def _split_text(cls, text: str, separators: List[str]) -> List[str]:
        """根据分隔符列表分割文本"""
        if not separators:
            return [text]
        
        separator = separators[0]
        parts = text.split(separator)
        
        # 如果只有一个部分或没有更多分隔符，返回结果
        if len(parts) == 1 or len(separators) == 1:
            return [part for part in parts if part.strip()]
        
        # 递归分割剩余部分
        result = []
        for part in parts:
            if part.strip():
                sub_parts = cls._split_text(part, separators[1:])
                result.extend(sub_parts)
        
        return result

    @staticmethod
    def _get_cover_content(text: str, cover_length: int) -> str:
        """获取文本末尾的覆盖内容"""
        if len(text) <= cover_length:
            return text
        
        # 尝试找到合适的分割点
        separators = ['\n', '. ', '! ', '? ', '; ', ', ', ' ']
        for sep in separators:
            pos = text.rfind(sep, len(text)-cover_length*2, len(text))
            if pos != -1:
                return text[pos+len(sep):]
        
        # 如果找不到分隔符，直接截取
        return text[-cover_length:]

    @classmethod
    def read_file_content(cls, file_path: str):
        # 根据文件扩展名选择读取方法
        if file_path.endswith('.pdf'):
            return cls.read_pdf(file_path)
        elif file_path.endswith('.md'):
            return cls.read_markdown(file_path)
        elif file_path.endswith('.txt'):
            return cls.read_text(file_path)
        else:
            raise ValueError("Unsupported file type")

    @classmethod
    def read_pdf(cls, file_path: str):
        """使用pdfplumber读取PDF文件，提高文本提取准确率"""
        try:
            with pdfplumber.open(file_path) as pdf:
                text = ""
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + '\n'
                return text.strip()
        except Exception as e:
            logger.error(f"读取PDF文件失败: {file_path}, 错误: {str(e)}")
            raise

    @classmethod
    def read_docx(cls, file_path: str):
        """读取DOCX文件内容"""
        try:
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            return text
        except Exception as e:
            logger.error(f"读取DOCX文件失败: {file_path}, 错误: {str(e)}")
            raise

    @classmethod
    def read_pptx(cls, file_path: str):
        """读取PPTX文件内容"""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text.append(shape.text)
            return "\n\n".join(text)
        except Exception as e:
            logger.error(f"读取PPTX文件失败: {file_path}, 错误: {str(e)}")
            raise

    @classmethod
    def read_markdown(cls, file_path: str):
        # 读取Markdown文件
        with open(file_path, 'r', encoding='utf-8') as file:
            md_text = file.read()
            html_text = markdown.markdown(md_text)
            # 使用BeautifulSoup从HTML中提取纯文本
            soup = BeautifulSoup(html_text, 'html.parser')
            plain_text = soup.get_text()
            # 使用正则表达式移除网址链接
            text = re.sub(r'http\S+', '', plain_text) 
            return text

    @classmethod
    def read_text(cls, file_path: str):
        # 读取文本文件
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()


class Documents:
    """
        获取已分好类的json格式文档
    """
    def __init__(self, path: str = '') -> None:
        self.path = path
    
    def get_content(self):
        with open(self.path, mode='r', encoding='utf-8') as f:
            content = json.load(f)
        return content